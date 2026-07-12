from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # 10 Slide definitions
    slides_content = [
        {"title": "Esillio", "subtitle": "Your body remembers. Now your healthcare can too."},
        {"title": "Market Opportunity", "text": "Trillions at stake. Digital Health projected to reach $2.3T by 2034."},
        {"title": "The Problem", "text": "Fragmented data siloes hinder care and innovation."},
        {"title": "The Solution", "text": "Biological Continuity Compiler™: A single, continuous health narrative."},
        {"title": "Technology", "text": "Open-source LLMs + AMD-accelerated edge inference."},
        {"title": "Business Model", "text": "Freemium for consumers; B2B licensing for clinics."},
        {"title": "Competitive Moat", "text": "Local-first privacy. Absolute data ownership."},
        {"title": "Roadmap", "text": "2026: MVP. 2027: Beta. 2028: Scale & API Launch."},
        {"title": "Why Invest?", "text": "Trillion-dollar wave. Proprietary continuity IP."},
        {"title": "Vision", "text": "Join us in building the OS your body has always needed."}
    ]

    for slide_data in slides_content:
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
        
        # Set matte black background (neutral-background)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(10, 10, 10)

        # Create Card (neutral-surface)
        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        height = Inches(5.5)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(13, 13, 13)
        shape.line.color.rgb = RGBColor(26, 26, 26) # border-primary (rgba 0.1)
        shape.line.width = Pt(1.5)

        # Title
        title = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(1))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data['title']
        p.font.size = Pt(44)
        p.font.color.rgb = RGBColor(255, 69, 51) # brand-primary
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Body text / Subtitle
        body = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(3))
        tf = body.text_frame
        p = tf.paragraphs[0]
        
        # Use subtitle if it exists, otherwise use text
        p.text = slide_data.get('subtitle', slide_data.get('text', ''))
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(150, 150, 150) # text-secondary
        p.alignment = PP_ALIGN.CENTER

    prs.save('Esillio_Pitch_Deck.pptx')
    print("Presentation created successfully: Esillio_Pitch_Deck.pptx")

if __name__ == "__main__":
    create_presentation()
