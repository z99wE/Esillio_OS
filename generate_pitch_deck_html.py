from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def apply_text_formatting(run, font_name, size_pt, rgb, bold=False, italic=False):
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.color.rgb = rgb
    run.font.bold = bold
    run.font.italic = italic

def create_presentation():
    prs = Presentation()
    
    # Colors from the design system
    bg_color = RGBColor(5, 5, 5) # Deep black
    text_main = RGBColor(255, 255, 255)
    text_muted = RGBColor(209, 213, 219)
    brand_orange = RGBColor(255, 149, 0)
    brand_pink = RGBColor(213, 73, 213)
    brand_violet = RGBColor(106, 13, 173)
    
    slides_data = [
        {
            "layout": "title",
            "title": "Esillio",
            "subtitle": "Your body remembers. Now your healthcare can too.\n\n\"Biology is memory; Esillio is its librarian.\"\n\nA local-first AI operating system that compiles your scattered health records into a secure, searchable, continuous timeline."
        },
        {
            "layout": "content",
            "title": "Trillions of dollars are at stake.",
            "subtitle": "\"The data exists – it just needs to be connected.\"",
            "points": [
                "$2.3T Digital Health: Growing aggressively from ~$450B today (projected 2033), demanding entirely new infrastructure.",
                "$195B Healthcare AI: Surging from ~$37B (2026), expected by 2031 to completely redefine clinical diagnostic tools.",
                "$64.2B Personal Apps: Consumer wellness app spending by 2033, proving the massive consumer appetite for data ownership."
            ]
        },
        {
            "layout": "content",
            "title": "No single view of your health story.",
            "subtitle": "\"Fragmented data costs lives and stalls innovation.\"",
            "points": [
                "Violently Siloed: Today’s healthcare data is scattered across disparate labs, EHRs, wearables, and illegible doctor notes. Fragmentation hinders care and innovation.",
                "Missing Context: Patients lack context: preventable errors and misdiagnoses rise dramatically when providers only see scattered pieces of the puzzle.",
                "Walled Gardens: Nearly all health data gets locked away in institutional silos, leaving individuals entirely unable to leverage their own history."
            ]
        },
        {
            "layout": "content",
            "title": "We compile your scattered history into one continuous narrative.",
            "subtitle": "\"The best time to organize your medical history was years ago. The second best time is today.\"",
            "points": [
                "Biological Compiler™: AI transforms any input (doctor notes, complex PDFs, CSVs, wearables) into structured, time‑stamped HealthEvents effortlessly.",
                "Persistent Health Graph: Links biomarkers, meds, and symptoms on-device. Enables trend alerts and complex semantic search.",
                "Local-First Privacy: Zero cloud dependency. This is your most sensitive data, governed entirely by your privacy, running securely directly on your machine."
            ]
        },
        {
            "layout": "content",
            "title": "AI innovation, optimized for privacy.",
            "subtitle": "\"Intelligence requires context, and context requires power.\"",
            "points": [
                "AMD-Accelerated AI Edge: Under the hood is an open-source LLM fine-tuned specifically on healthcare data. We heavily leverage AMD’s high-performance compute (ROCm) for development, enabling hyper-efficient on-device inference.",
                "Vector + SQL Hybrid: ChromaDB and heavily encrypted SQLite hold the timeline. This architecture allows the OS to answer highly complex free-text health queries instantly, offline. Future-proofed for genomics, voice, and visual inputs."
            ]
        },
        {
            "layout": "content",
            "title": "Accessible to everyone, with enterprise optionality.",
            "subtitle": "\"A sustainable ecosystem built on undeniable value.\"",
            "points": [
                "Freemium App: Basic timeline tracking is free. Premium features via subscription ($5–$10/mo).",
                "B2B Licensing (Core Revenue): Clinics, insurers, and pharma license our HIPAA-compliant API to drastically improve chronic care management and patient outreach.",
                "$100M+ Projection: Capturing just 0.1% of the digital health TAM yields $4M+ ARR. Scaling to $50–$100M+ revenue is highly feasible by 2028."
            ]
        },
        {
            "layout": "content",
            "title": "No one else connects the dots.",
            "subtitle": "\"True health memory requires absolute data sovereignty.\"",
            "points": [
                "Traditional Big Tech: Rely on cloud processing, inherently risking massive data privacy breaches. Fragmented entry forces users into isolated silos with zero longitudinal memory.",
                "Esillio OS: HIPAA-grade, entirely on-device processing. Your data never leaves. We create a lifelong narrative. Once data is in Esillio, lock-in is immense. Optimized local AMD infrastructure prevents open-source copycats."
            ]
        },
        {
            "layout": "content",
            "title": "Building the future of personal health.",
            "subtitle": "\"A relentless march toward proactive healthcare.\"",
            "points": [
                "2026 (Done): MVP Continuity Compiler complete. Manual logging and basic HealthEvent model successfully prototyped in app UI. Early validation.",
                "2027: Add complex AI ingestion (OCR/PDF parsing, NLP for clinical notes, wearables). Launch closed beta pilots with healthcare partners.",
                "2028: Public multi-platform app launch. Developer API (Vector DB) release. First revenue generation backed by AMD-accelerated model training.",
                "2029+: Expand globally, add specialized agents, and execute continuous model improvements."
            ]
        },
        {
            "layout": "content",
            "title": "We're the first to build a true health OS.",
            "subtitle": "First-Mover Advantage: No major player does this with absolute privacy and local AI. Our IP is novel, defensible, and scales globally.",
            "points": [
                "$2T+ Market Wave: Even niche success natively taps into a trillion-dollar TAM.",
                "100% Local Processing: Absolute privacy guarantee via secure on-device AI.",
                "MVP Validated: Hackathon MVP complete. Early local-first user tests are highly positive.",
                "The Ask: Raising Seed Capital to finalize complex AI integrations, optimize our AMD inference pipeline, and launch clinical beta pilots."
            ]
        },
        {
            "layout": "title",
            "title": "Esillio",
            "subtitle": "Your health history is your most valuable data.\nWe give it a voice.\n\n\"Join us in building the OS that your body has always needed.\"\n\nInvest in Esillio (founders@esillio.com)"
        }
    ]

    for index, slide_data in enumerate(slides_data):
        slide_layout = prs.slide_layouts[6] # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        if slide_data["layout"] == "title":
            # Title
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = slide_data["title"]
            apply_text_formatting(run, "Playfair Display", 64, text_main, bold=True)
            
            # Subtitle
            txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(2))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run()
            run2.text = slide_data["subtitle"]
            apply_text_formatting(run2, "Inter", 18, text_muted)
        
        else:
            # Title
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = slide_data["title"]
            apply_text_formatting(run, "Playfair Display", 40, text_main, bold=True)
            
            # Subtitle
            if slide_data.get("subtitle"):
                txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.5))
                tf2 = txBox2.text_frame
                tf2.word_wrap = True
                p2 = tf2.paragraphs[0]
                run2 = p2.add_run()
                run2.text = slide_data["subtitle"]
                apply_text_formatting(run2, "Playfair Display", 20, brand_pink, italic=True)
                
            # Content
            if slide_data.get("points"):
                txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(4.5))
                tf3 = txBox3.text_frame
                tf3.word_wrap = True
                
                for pt in slide_data["points"]:
                    p3 = tf3.add_paragraph()
                    p3.level = 0
                    run3 = p3.add_run()
                    run3.text = "• " + pt
                    apply_text_formatting(run3, "Inter", 18, text_muted)
                    
                    # Add space after
                    p3.space_after = Pt(14)
                    
    prs.save('Esillio_Pitch_Deck_From_HTML.pptx')

if __name__ == '__main__':
    create_presentation()
